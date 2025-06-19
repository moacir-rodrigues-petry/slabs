#!/usr/bin/env python
"""
Script to convert PyChat_Presentation.md to a PowerPoint presentation.
This script uses python-pptx to create a PPTX file based on the markdown content.
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle=None, presenter=None, date=None):
    """Create a title slide with optional subtitle, presenter, and date."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(44, 77, 121)  # Dark blue
    
    # Subtitle
    subtitle_shape = slide.placeholders[1]
    tf = subtitle_shape.text_frame
    tf.clear()
    
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(32)
        p.font.bold = False
        p.font.italic = True
        
    if presenter:
        p = tf.add_paragraph()
        p.text = ""  # Add a blank line
        
        p = tf.add_paragraph()
        p.text = presenter
        p.font.size = Pt(24)
        
    if date:
        p = tf.add_paragraph()
        p.text = ""  # Add a blank line
        
        p = tf.add_paragraph()
        p.text = date
        p.font.size = Pt(20)
        p.font.italic = True
    
    return slide

def create_content_slide(prs, title, content):
    """Create a content slide with bullet points."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Content
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.clear()  # Clear existing content
    
    # Process content line by line for bullet points
    lines = content.split('\n')
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        p = tf.add_paragraph()
        
        # Check for bullet points
        if line.startswith('- ') or line.startswith('• '):
            p.text = line[2:].strip()
            p.level = 0
        elif line.startswith('  - ') or line.startswith('  • '):
            p.text = line[4:].strip()
            p.level = 1
        elif line.startswith('    - ') or line.startswith('    • '):
            p.text = line[6:].strip()
            p.level = 2
        else:
            p.text = line
            p.level = 0
            
        # Process checkmarks and emoji indicators
        text = p.text
        if '✅' in text:
            text = text.replace('✅', '')
            p.text = text.strip()
            p.font.color.rgb = RGBColor(0, 128, 0)  # Green for completed items
        elif '🔄' in text:
            text = text.replace('🔄', '')
            p.text = text.strip()
            p.font.color.rgb = RGBColor(255, 140, 0)  # Orange for in-progress items
        elif '⏳' in text:
            text = text.replace('⏳', '')
            p.text = text.strip()
            p.font.color.rgb = RGBColor(128, 128, 128)  # Gray for planned items
            
    return slide

def process_markdown_content(md_content):
    """Process markdown content and extract slides."""
    slides = []
    slide_sections = re.split(r'---\s*\n', md_content)
    
    for section in slide_sections:
        section = section.strip()
        if not section:
            continue
        
        lines = section.split('\n')
        slide_info = {"title": "", "content": [], "is_title_slide": False}
        
        # Extract slide title (first h1 or h2)
        for i, line in enumerate(lines):
            if line.startswith('## Slide'):
                continue
            if line.startswith('# '):
                slide_info["title"] = line[2:].strip()
                lines = lines[i+1:]
                break
        
        # Determine if it's the title slide
        if "Slide 1:" in section:
            slide_info["is_title_slide"] = True
            
            # Extract date if present
            date_match = re.search(r'_([^_]+)_\s*$', section, re.MULTILINE)
            if date_match:
                slide_info["date"] = date_match.group(1).strip()
        
        # Process content (ignoring mermaid diagrams for now)
        content_text = ""
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            
            # Skip mermaid blocks
            if line.startswith('```mermaid'):
                while i < len(lines) and not lines[i].strip().endswith('```'):
                    i += 1
                if i < len(lines):  # Skip the closing ```
                    i += 1
                continue
                
            # Skip other code blocks
            elif line.startswith('```'):
                while i < len(lines) and not lines[i].strip().endswith('```'):
                    i += 1
                if i < len(lines):  # Skip the closing ```
                    i += 1
                continue
            
            # Skip image links
            elif line.startswith('!['):
                i += 1
                continue
                
            # Process headings 
            elif line.startswith('## ') or line.startswith('### '):
                content_text += line.replace('##', '').replace('#', '').strip() + "\n\n"
                i += 1
                continue
            
            # Process bullet points and status indicators
            elif line.startswith('-') or line.startswith('•'):
                # Keep bullet points with their status indicators (✅, 🔄, ⏳)
                content_text += line + "\n"
                i += 1
                continue
                
            # Add regular content
            elif line and not line.startswith('##') and not line.startswith('#'):
                content_text += line + "\n"
            
            i += 1
        
        slide_info["content"] = content_text.strip()
        slides.append(slide_info)
    
    return slides

def create_demo_slide(prs, title, image_path=None, content=None):
    """Create a demo slide with an image and optional content."""
    slide_layout = prs.slide_layouts[6]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Add content if provided
    if content:
        content_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
        content_frame = content_box.text_frame
        p = content_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER
    
    # If we had an actual image, we would add it here
    # For now, we'll just add a placeholder text box
    if not image_path or not os.path.exists(image_path):
        img_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(3))
        img_frame = img_box.text_frame
        p = img_frame.add_paragraph()
        p.text = "[PyChat Demo Screenshot]"
        p.font.size = Pt(32)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
    else:
        slide.shapes.add_picture(image_path, Inches(2), Inches(2), Inches(9), Inches(3))
    
    return slide

def create_code_slide(prs, title, code_content):
    """Create a slide with code block."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    
    # Code content - use a textbox with monospace font
    code_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
    tf = code_box.text_frame
    tf.word_wrap = True
    
    # Split code by lines and add as paragraphs
    lines = code_content.split('\n')
    for i, line in enumerate(lines):
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Courier New"
        p.font.size = Pt(16)
    
    return slide

def create_ppt_from_markdown(md_file, output_file):
    """Create a PowerPoint presentation from a markdown file."""
    # Read markdown content
    with open(md_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    # Process markdown content
    slides = process_markdown_content(md_content)
    
    # Create presentation
    prs = Presentation()
    
    # Define slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    for i, slide_info in enumerate(slides):
        title = slide_info["title"]
        content = slide_info["content"]
        
        # Handle different slide types
        if i == 0 or slide_info.get("is_title_slide", False):
            # Title slide
            date = slide_info.get("date", "June 19, 2025")
            create_title_slide(prs, title, subtitle="A Simple Python Chat Application", date=date)
        
        # Special handling for code examples slide
        elif "Running the Application" in title or "bash" in content:
            # Extract code blocks if present
            code_content = ""
            for line in content.split('\n'):
                if line.strip() and not line.startswith('-'):
                    code_content += line + "\n"
            create_code_slide(prs, title, code_content)
        
        # Special handling for Demo slide
        elif "Demo & Questions" in title:
            create_demo_slide(prs, title, content="Thank you for your attention!\n\nQuestions?")
        
        # Regular content slides
        else:
            create_content_slide(prs, title, content)
    
    # Save presentation
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    try:
        script_dir = os.path.dirname(os.path.abspath(__file__))
        md_file = os.path.join(script_dir, "PyChat_Presentation.md")
        output_file = os.path.join(script_dir, "PyChat_Presentation.pptx")
        
        print(f"Reading markdown from: {md_file}")
        print(f"Will save presentation to: {output_file}")
        
        if not os.path.exists(md_file):
            print(f"Error: Markdown file not found at {md_file}")
            exit(1)
            
        create_ppt_from_markdown(md_file, output_file)
    except Exception as e:
        import traceback
        print(f"Error: {str(e)}")
        traceback.print_exc()
